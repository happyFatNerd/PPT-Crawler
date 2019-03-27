from pptx import Presentation
import re

def crawlPpt(url,keywords):
    try:
        prs = Presentation(url)
    except:
        return 0
    content = []
    for slide in prs.slides:
        text = ''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue        
            text += shape.text
            '''for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)'''
        try:
            info = []
            if(keywords[0]==''):
                return -1
            splitText(text,len(keywords),info,keywords)
            content.append(info)
        except:
            continue
    return content

def splitText(text,num,info,keywords):
    texts = text.split(keywords[len(keywords)-num])
    text = texts[1]
    info.append(texts[0])
    num-=1
    if num>0:
        splitText(text,num,info,keywords)
    if num==0:
        info.append(text)

def printText(content,keywords):
    fullText = ''
    if content == None:
        print('Don\'t have such keywords in one single page!')
        input()
    else:
        for info in content:
            for i in range(len(keywords)+1):
                if i ==0:
                    fullText += 'Title:\n'+re.sub('\n+','\n',info[i])+'\n\n'
                else:
                    fullText += keywords[i-1]+':\n'+re.sub('\n+','\n',info[i])+'\n\n'
            fullText += '____________________________________________\n\n'
        return fullText

def main():
    url = input('Please input the file path(including suffix):')
    url = url.replace('\\','\\\\')
    num = 1
    keywords = []
    try:
        num = int(input('How many keywords you need:'))
    except:
        print('Invaild number! Default as 1 keyword.')
    finally:
        for i in range(num):
            keywords.append(input('Input keyword'+str(i+1)+':'))
    if crawlPpt(url,keywords) == 0:
        print('url not existing!')
        input()
        return
    if crawlPpt(url,keywords) == -1:
        print('no keyword!')
        input()
        return
    fullText = printText(crawlPpt(url,keywords),keywords)
    with open(url.split('\\\\')[-1].split('.')[0]+'_text.txt','w') as f:
        f.write(fullText)

main()
